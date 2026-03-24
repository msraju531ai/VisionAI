"""Generate VisionAI Demo PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x0D, 0x0D)   # near-black
PANEL     = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
ACCENT    = RGBColor(0x19, 0xC3, 0x7D)   # green
RED       = RGBColor(0xEF, 0x44, 0x44)
YELLOW    = RGBColor(0xF5, 0x9E, 0x0B)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x11, 0x11, 0x11)

SW, SH = Inches(13.33), Inches(7.5)   # 16:9 widescreen


def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    from pptx.util import Pt as Pt2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt2(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt2(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def add_bullet_box(slide, bullets, l, t, w, h,
                   font_size=16, color=LIGHT, bold_first=False,
                   dot_color=ACCENT):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold  = bold_first and (item == bullets[0])
    return txb


def section_header_bar(slide, label, color=ACCENT):
    """Thin coloured bar at top with section label."""
    add_rect(slide, 0, 0, 13.33, 0.55, color)
    add_textbox(slide, label, 0.3, 0.06, 12, 0.45,
                font_size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  — Title
# ══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)

    # diagonal accent strip
    add_rect(s, 0, 0, 13.33, 0.12, ACCENT)
    add_rect(s, 0, 7.38, 13.33, 0.12, ACCENT)

    # left accent bar
    add_rect(s, 0, 0, 0.18, 7.5, PANEL)

    # centre card
    add_rect(s, 1.5, 1.4, 10.33, 4.7, PANEL)

    # logo text
    add_textbox(s, "VisionAI", 1.8, 1.7, 6, 1.0,
                font_size=52, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

    add_textbox(s, "AI-Powered Workforce Intelligence", 1.8, 2.65, 10, 0.7,
                font_size=26, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(s, "Demo Presentation", 1.8, 3.35, 10, 0.55,
                font_size=20, bold=False, color=LIGHT, align=PP_ALIGN.LEFT, italic=True)

    # three coloured pills representing 3 demos
    pills = [
        ("Demo 1  Employee Detection", ACCENT),
        ("Demo 2  Unauthorized Entry Detection", RED),
        ("Demo 3  Idle + Shift Compliance", YELLOW),
    ]
    for i, (label, col) in enumerate(pills):
        x = 1.8 + i * 3.5
        add_rect(s, x, 4.35, 3.2, 0.45, col)
        add_textbox(s, label, x + 0.1, 4.36, 3.1, 0.43,
                    font_size=11, bold=True, color=DARK_TEXT, align=PP_ALIGN.LEFT)

    add_textbox(s, "Blue Waters Project  |  2026", 1.8, 5.1, 10, 0.4,
                font_size=14, color=LIGHT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  — Agenda
# ══════════════════════════════════════════════════════════════════════════════
def slide_agenda(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Agenda", BLUE)

    items = [
        ("01", "Demo 1 — Employee Detection",            ACCENT,
         "Register employees  ›  Upload video  ›  View annotated output with ID labels"),
        ("02", "Demo 2 — Unauthorized Entry Detection", RED,
         "Face matching against employee DB  ›  Authorized (Green) vs Unauthorized (Red)  ›  Event log"),
        ("03", "Demo 3A — Idle Time Detection",           YELLOW,
         "Track person movement  ›  Yellow box when still  ›  Idle event log"),
        ("04", "Demo 3B — Shift Compliance",              BLUE,
         "Map video timestamp to wall clock  ›  Compare vs. work schedule  ›  Compliance report"),
    ]

    for i, (num, title, col, desc) in enumerate(items):
        y = 0.85 + i * 1.55
        add_rect(s, 0.4, y, 0.7, 0.7, col)
        add_textbox(s, num, 0.4, y, 0.7, 0.7,
                    font_size=22, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(s, title, 1.25, y - 0.04, 11.5, 0.5,
                    font_size=18, bold=True, color=WHITE)
        add_textbox(s, desc, 1.25, y + 0.44, 11.5, 0.6,
                    font_size=13, color=LIGHT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  — Demo 1 Overview
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo1_overview(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 1  —  Employee Detection", ACCENT)

    add_textbox(s, "What it does", 0.5, 0.75, 12, 0.5,
                font_size=20, bold=True, color=ACCENT)
    add_textbox(s,
        "Automatically identifies registered employees in an uploaded video using "
        "deep-learning face embeddings (Facenet) and annotates each person with their name and confidence score.",
        0.5, 1.2, 12.3, 0.8, font_size=15, color=LIGHT)

    # two columns
    # LEFT — How it works
    add_rect(s, 0.4, 2.15, 5.8, 4.6, PANEL)
    add_textbox(s, "How It Works", 0.65, 2.25, 5.3, 0.45,
                font_size=16, bold=True, color=ACCENT)
    steps = [
        "1  Upload employee photo via /demo page",
        "2  Facenet extracts 128-dim face embedding",
        "3  Embedding saved in SQLite DB",
        "4  Upload a workplace video",
        "5  YOLO v8 detects every person per frame",
        "6  Face crop → Facenet embedding → cosine similarity match",
        "7  Green box + name label for recognized employees",
        "8  Red box for unrecognized persons",
    ]
    txb = s.shapes.add_textbox(Inches(0.65), Inches(2.75), Inches(5.3), Inches(3.8))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = step
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT

    # RIGHT — Output legend
    add_rect(s, 6.7, 2.15, 6.2, 4.6, PANEL)
    add_textbox(s, "Output Video Legend", 6.95, 2.25, 5.7, 0.45,
                font_size=16, bold=True, color=ACCENT)
    legend = [
        (ACCENT, "Green box  — Recognized employee  (name + score shown)"),
        (RED,    "Red box    — Unknown / Unregistered person"),
    ]
    for j, (col, label) in enumerate(legend):
        y = 2.85 + j * 0.8
        add_rect(s, 7.0, y, 0.35, 0.35, col)
        add_textbox(s, label, 7.5, y - 0.03, 5.1, 0.5,
                    font_size=13, color=LIGHT)

    add_textbox(s, "Key Technology", 6.95, 4.6, 5.5, 0.4,
                font_size=14, bold=True, color=ACCENT)
    techs = ["• YOLOv8 — person bounding box detection",
             "• DeepFace / Facenet — face embedding (128-dim)",
             "• Cosine similarity — identity matching",
             "• ffmpeg — H.264 annotated video output"]
    for k, tech in enumerate(techs):
        add_textbox(s, tech, 6.95, 5.05 + k * 0.38, 5.5, 0.4,
                    font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  — Demo 1 How to Test
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo1_test(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 1  —  How to Test", ACCENT)

    steps = [
        ("Step 1", "Register Employees",
         "Go to /demo  →  Employee Registration section\n"
         "Enter name, upload a clear front-facing photo  →  Click Register"),
        ("Step 2", "Upload Video",
         "Go to /demo  →  Uploaded Videos section\n"
         "Choose a video file containing the registered employee  →  Click Upload & Process"),
        ("Step 3", "Watch Progress",
         "The Processed Frames counter updates in real-time (no page refresh needed)\n"
         "Status changes: uploaded  →  processing  →  completed"),
        ("Step 4", "View Annotated Output",
         "Click the View button in the Processed Videos table\n"
         "Green box = identified employee with name label\n"
         "Red box = unknown person"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 0.8 + i * 1.55
        add_rect(s, 0.4, y, 1.5, 0.55, ACCENT)
        add_textbox(s, num, 0.4, y, 1.5, 0.55,
                    font_size=13, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(s, title, 2.1, y, 10.8, 0.5,
                    font_size=16, bold=True, color=WHITE)
        add_textbox(s, desc, 2.1, y + 0.5, 10.8, 0.9,
                    font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  — Demo 2 Overview
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo2_overview(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 2  —  Unauthorized Entry Detection", RED)

    add_textbox(s, "What it does", 0.5, 0.75, 12, 0.5,
                font_size=20, bold=True, color=RED)
    add_textbox(s,
        "Scans an uploaded video and compares every detected person's face against the "
        "registered employee database. Each person is labelled Authorized (Green) or "
        "Unauthorized (Red) in the annotated output video, and unauthorized events are logged.",
        0.5, 1.2, 12.3, 0.9, font_size=15, color=LIGHT)

    # flow boxes  — no email step
    flow = [
        ("Upload Video", PANEL),
        ("YOLO v8\nDetects Persons", PANEL),
        ("Facenet\nEmbedding", PANEL),
        ("Cosine Match\nvs Employee DB", PANEL),
        ("Authorized\n(Green)", ACCENT),
        ("Unauthorized\n(Red)", RED),
    ]
    arrow = "→"
    n = len(flow)
    box_w = 1.85
    gap   = 0.22
    total = n * box_w + (n - 1) * gap
    start_x = (13.33 - total) / 2

    for i, (label, col) in enumerate(flow):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, 2.25, box_w, 0.95, col)
        txt_color = DARK_TEXT if col in (ACCENT, RED) else WHITE
        add_textbox(s, label, x, 2.25, box_w, 0.95,
                    font_size=12, bold=True, color=txt_color,
                    align=PP_ALIGN.CENTER)
        if i < n - 1:
            add_textbox(s, arrow, x + box_w, 2.5, gap + 0.05, 0.5,
                        font_size=16, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)

    # ── THREE detail cards ──────────────────────────────────────────────────
    # Card 1 — Authorized Logic
    add_rect(s, 0.4, 3.45, 4.0, 3.8, PANEL)
    add_rect(s, 0.4, 3.45, 0.15, 3.8, ACCENT)
    add_textbox(s, "Authorized Person", 0.7, 3.55, 3.6, 0.48,
                font_size=14, bold=True, color=ACCENT)
    auth_pts = [
        "• Face detected in video frame",
        "• Facenet 128-dim embedding extracted",
        "• Cosine similarity >= 0.55 against a registered employee embedding",
        "• GREEN bounding box drawn",
        "• Employee name + confidence score shown as label",
        "• Person's identity confirmed and logged",
    ]
    txb = s.shapes.add_textbox(Inches(0.7), Inches(4.1), Inches(3.55), Inches(3.0))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(auth_pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = LIGHT

    # Card 2 — Unauthorized Logic
    add_rect(s, 4.67, 3.45, 4.0, 3.8, PANEL)
    add_rect(s, 4.67, 3.45, 0.15, 3.8, RED)
    add_textbox(s, "Unauthorized Person", 4.97, 3.55, 3.6, 0.48,
                font_size=14, bold=True, color=RED)
    unauth_pts = [
        "• Face detected but similarity < 0.55 for ALL employees",
        "• OR no face detected in person crop",
        "• RED bounding box drawn",
        "• Label shows 'Unauthorized'",
        "• Unauthorized entry event recorded in DB",
        "• Timestamp and video ID stored for audit",
    ]
    txb2 = s.shapes.add_textbox(Inches(4.97), Inches(4.1), Inches(3.55), Inches(3.0))
    txb2.word_wrap = True
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(unauth_pts):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = LIGHT

    # Card 3 — Key Technology
    add_rect(s, 8.93, 3.45, 4.0, 3.8, PANEL)
    add_rect(s, 8.93, 3.45, 0.15, 3.8, BLUE)
    add_textbox(s, "Key Technology", 9.23, 3.55, 3.6, 0.48,
                font_size=14, bold=True, color=BLUE)
    tech_pts = [
        "• YOLOv8 — person bounding box (COCO class 0)",
        "• DeepFace / Facenet — 128-dim face embedding",
        "• Cosine similarity — identity matching",
        "• Similarity threshold: 0.55",
        "• last_boxes pattern — stable annotations on every frame",
        "• ffmpeg pipe — H.264 browser-playable output",
    ]
    txb3 = s.shapes.add_textbox(Inches(9.23), Inches(4.1), Inches(3.55), Inches(3.0))
    txb3.word_wrap = True
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    for i, line in enumerate(tech_pts):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = LIGHT


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  — Demo 2 How to Test
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo2_test(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 2  —  How to Test", RED)

    steps = [
        ("Step 1", "Ensure Employees Are Registered",
         "Go to /demo and confirm your employees appear in the Registered Employees table\n"
         "If empty — register them with a clear front-facing photo first"),
        ("Step 2", "Upload Video on /demo/unauthorized",
         "Go to /demo/unauthorized  →  Upload a video that contains both a registered employee AND a stranger\n"
         "Click Upload & Process"),
        ("Step 3", "Watch Live Progress",
         "Progress bar and Processed Frames counter update every 2 seconds — no page refresh needed\n"
         "Status changes: uploaded  →  processing  →  completed"),
        ("Step 4", "Review Annotated Output",
         "Click View when status = completed\n"
         "Green box + name = authorized employee  |  Red box = unauthorized stranger\n"
         "Unauthorized Entry Events table below lists every flagged person with timestamp"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = 0.8 + i * 1.55
        add_rect(s, 0.4, y, 1.5, 0.55, RED)
        add_textbox(s, num, 0.4, y, 1.5, 0.55,
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, title, 2.1, y, 10.8, 0.5,
                    font_size=16, bold=True, color=WHITE)
        add_textbox(s, desc, 2.1, y + 0.5, 10.8, 0.9,
                    font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  — Demo 3A — Idle Detection
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo3a(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 3A  —  Idle Time Detection", YELLOW)

    add_textbox(s, "What it does", 0.5, 0.75, 12, 0.5,
                font_size=20, bold=True, color=YELLOW)
    add_textbox(s,
        "Tracks each detected person's centroid across video frames. "
        "When a person has not moved more than 30 pixels for longer than the idle threshold, "
        "their bounding box turns yellow and an idle event (start time, end time, duration) is recorded.",
        0.5, 1.2, 12.3, 0.85, font_size=15, color=LIGHT)

    add_rect(s, 0.4, 2.2, 12.5, 0.06, YELLOW)

    # Three info cards
    cards = [
        ("Movement Tracking", YELLOW,
         ["• Centroid (centre-point) tracked per person per frame",
          "• Movement threshold: 30 pixels",
          "• Idle clock starts when movement < threshold",
          "• Idle clock resets on any significant movement"]),
        ("Idle Threshold", YELLOW,
         ["• Configurable in config/settings.py",
          "• Default: 300 seconds (5 minutes)",
          "• For demo: lower to 10 seconds",
          "• idle_threshold_seconds = 10"]),
        ("Output", YELLOW,
         ["• Yellow box + 'Idle: <name>' label on video",
          "• Idle Events table: start / end / duration",
          "• Employee ID linked if person is recognized",
          "• All events persisted in SQLite"]),
    ]
    for i, (title, col, bullets) in enumerate(cards):
        x = 0.4 + i * 4.35
        add_rect(s, x, 2.45, 4.1, 4.8, PANEL)
        add_textbox(s, title, x + 0.2, 2.55, 3.7, 0.5,
                    font_size=15, bold=True, color=col)
        txb = s.shapes.add_textbox(Inches(x + 0.2), Inches(3.1), Inches(3.7), Inches(4.0))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        for j, b in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_before = Pt(7)
            run = p.add_run()
            run.text = b
            run.font.size = Pt(13)
            run.font.color.rgb = LIGHT

    add_textbox(s,
        "⚠  IMPORTANT for demo:  Lower idle_threshold_seconds to 10 in config/settings.py and restart server",
        0.4, 7.1, 12.5, 0.35, font_size=12, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  — Demo 3B — Shift Compliance
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo3b(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 3B  —  Shift Compliance (Working Hours)", BLUE)

    add_textbox(s, "What it does", 0.5, 0.75, 12, 0.5,
                font_size=20, bold=True, color=BLUE)
    add_textbox(s,
        "Maps video frame timestamps to real wall-clock time using the Video Start DateTime. "
        "Compares each employee's first appearance (check-in) and last appearance (check-out) "
        "against their configured work schedule, and generates a compliance report.",
        0.5, 1.2, 12.3, 0.9, font_size=15, color=LIGHT)

    # Status cards
    statuses = [
        ("compliant", ACCENT,  "Employee arrived and left within the grace period of their schedule"),
        ("late",      YELLOW,  "Employee's first appearance is after expected start + grace minutes"),
        ("early_exit",YELLOW,  "Employee's last appearance is before expected end − grace minutes"),
        ("absent",    RED,     "Employee was not detected at all in the video"),
    ]
    for i, (status, col, desc) in enumerate(statuses):
        y = 2.25 + i * 1.15
        add_rect(s, 0.4, y, 2.5, 0.65, col)
        add_textbox(s, status.upper(), 0.4, y, 2.5, 0.65,
                    font_size=15, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(s, desc, 3.15, y + 0.1, 10.0, 0.55,
                    font_size=14, color=LIGHT)

    add_rect(s, 0.4, 6.85, 0.06, 0.5, BLUE)  # decorative

    add_textbox(s, "Compliance Table Columns:", 0.6, 6.82, 12, 0.35,
                font_size=13, bold=True, color=BLUE)
    cols_text = "Employee  |  Check In  |  Check Out  |  Total Minutes  |  Status  |  Deviation (minutes)"
    add_textbox(s, cols_text, 0.6, 7.15, 12.5, 0.35,
                font_size=12, color=LIGHT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9  — Demo 3 How to Test
# ══════════════════════════════════════════════════════════════════════════════
def slide_demo3_test(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "Demo 3  —  How to Test (Idle + Shift Compliance)", YELLOW)

    steps = [
        ("Prep", "Lower Idle Threshold",
         "Edit config/settings.py  →  set  idle_threshold_seconds = 10\nRestart server  (uvicorn main:app --reload)"),
        ("Step 1", "Set Work Schedule",
         "Go to /demo/idle  →  Work Schedule section\n"
         "Select employee, enter Expected Start (e.g. 09:00), Expected End (18:00), Grace 10 min  →  Save Schedule"),
        ("Step 2", "Upload Video with Start Time",
         "In Video Start DateTime enter the simulated recording start time\n"
         "e.g. 2026-03-24 09:30:00  (30 min late — expect 'late' status)\n"
         "Upload a video where the employee is visible and stands still for >10 s"),
        ("Step 3", "Check Results",
         "Processed Videos table updates live (Job, Processed Frames)\n"
         "Click View → Yellow boxes on still persons, Green on moving employees\n"
         "Scroll down → Idle Events table + Shift Compliance Summary (late / compliant / etc.)"),
    ]
    col_map = {"Prep": BLUE, "Step 1": YELLOW, "Step 2": YELLOW, "Step 3": YELLOW}
    for i, (num, title, desc) in enumerate(steps):
        y = 0.8 + i * 1.55
        col = col_map.get(num, YELLOW)
        add_rect(s, 0.4, y, 1.5, 0.55, col)
        add_textbox(s, num, 0.4, y, 1.5, 0.55,
                    font_size=13, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(s, title, 2.1, y, 10.8, 0.5,
                    font_size=16, bold=True, color=WHITE)
        add_textbox(s, desc, 2.1, y + 0.5, 10.8, 0.9,
                    font_size=13, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10  — Architecture Summary
# ══════════════════════════════════════════════════════════════════════════════
def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    section_header_bar(s, "System Architecture", ACCENT)

    layers = [
        ("Frontend (Jinja2 Templates)",   PANEL, ACCENT,
         "/demo   /demo/unauthorized   /demo/idle   — dark-themed dashboard, live polling (2 s interval)"),
        ("API Layer (FastAPI)",            PANEL, BLUE,
         "REST endpoints for employee CRUD, video upload, YouTube download, status polling, video streaming"),
        ("Processing Services",           PANEL, YELLOW,
         "UnauthorizedDemoProcessor  |  IdleDemoProcessor  |  EmployeeIdentifier (DeepFace/Facenet)\n"
         "ActivityAnalyzer (centroid tracker)  |  PersonDetector (YOLOv8)  |  AlertService (SMTP email)"),
        ("Storage",                       PANEL, RED,
         "SQLite via async SQLAlchemy  |  Video files on disk  |  Employee photos on disk\n"
         "Tables: demo_employees, unauthorized_demo_videos, idle_demo_videos, idle_events, work_schedules"),
    ]

    for i, (title, bg, col, desc) in enumerate(layers):
        y = 0.75 + i * 1.6
        add_rect(s, 0.4, y, 12.5, 1.35, bg)
        add_rect(s, 0.4, y, 0.18, 1.35, col)
        add_textbox(s, title, 0.75, y + 0.05, 12.0, 0.45,
                    font_size=15, bold=True, color=col)
        add_textbox(s, desc, 0.75, y + 0.52, 12.0, 0.8,
                    font_size=12, color=LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11  — Closing
# ══════════════════════════════════════════════════════════════════════════════
def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s, BG)
    add_rect(s, 0, 0, 13.33, 0.12, ACCENT)
    add_rect(s, 0, 7.38, 13.33, 0.12, ACCENT)

    add_textbox(s, "VisionAI", 0.5, 1.5, 12, 1.2,
                font_size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, "AI-Powered Workforce Intelligence", 0.5, 2.7, 12, 0.7,
                font_size=24, color=WHITE, align=PP_ALIGN.CENTER)

    summary = [
        ("Demo 1", "Employee Detection",              ACCENT),
        ("Demo 2", "Unauthorized Entry Detection",     RED),
        ("Demo 3", "Idle + Shift Compliance",          YELLOW),
    ]
    for i, (d, t, c) in enumerate(summary):
        x = 1.5 + i * 3.6
        add_rect(s, x, 3.7, 3.1, 0.55, c)
        add_textbox(s, f"{d}: {t}", x + 0.1, 3.71, 3.0, 0.53,
                    font_size=11, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

    add_textbox(s, "Thank You", 0.5, 4.7, 12, 0.9,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, "Blue Waters Project  |  2026", 0.5, 5.65, 12, 0.5,
                font_size=16, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = new_prs()
    slide_title(prs)
    slide_agenda(prs)
    slide_demo1_overview(prs)
    slide_demo1_test(prs)
    slide_demo2_overview(prs)
    slide_demo2_test(prs)
    slide_demo3a(prs)
    slide_demo3b(prs)
    slide_demo3_test(prs)
    slide_architecture(prs)
    slide_closing(prs)

    out = "VisionAI_Demo_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()